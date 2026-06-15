from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_DIR = Path(__file__).resolve().parents[1]
REPORT_DIR = PROJECT_DIR / "report"
FIGURE_DIR = PROJECT_DIR / "outputs" / "figures"
PPTX_PATH = REPORT_DIR / "EEG脑信号偏好推荐实验_答辩PPT.pptx"
OUTLINE_PATH = REPORT_DIR / "答辩PPT大纲.md"


NAVY = RGBColor(28, 42, 61)
BLUE = RGBColor(43, 111, 176)
TEAL = RGBColor(42, 157, 143)
GRAY = RGBColor(96, 108, 120)
LIGHT = RGBColor(245, 247, 250)
WHITE = RGBColor(255, 255, 255)


def add_title(slide, title: str, subtitle: str = ""):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(8.4), Inches(0.48))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Arial"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(0.82), Inches(8.1), Inches(0.26))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.name = "Arial"
        sp.font.size = Pt(9.5)
        sp.font.color.rgb = GRAY
    slide.shapes.add_shape(1, Inches(0.55), Inches(1.12), Inches(1.05), Inches(0.04)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = BLUE
    slide.shapes[-1].line.fill.background()


def add_footer(slide, page: int):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.03), Inches(8.2), Inches(0.22))
    p = box.text_frame.paragraphs[0]
    p.text = f"基于脑信号推断用户偏好的深度学习推荐实验  |  {page}/5"
    p.font.name = "Arial"
    p.font.size = Pt(8.5)
    p.font.color.rgb = GRAY


def add_bullets(slide, items, x, y, w, h, font_size=16, color=NAVY):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Arial"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return box


def add_metric_card(slide, x, y, label, value, note=""):
    card = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(1.65), Inches(0.92))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = RGBColor(224, 229, 235)
    box = slide.shapes.add_textbox(Inches(x + 0.12), Inches(y + 0.12), Inches(1.4), Inches(0.32))
    p = box.text_frame.paragraphs[0]
    p.text = value
    p.font.name = "Arial"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLUE
    lab = slide.shapes.add_textbox(Inches(x + 0.12), Inches(y + 0.52), Inches(1.4), Inches(0.22))
    lp = lab.text_frame.paragraphs[0]
    lp.text = label
    lp.font.name = "Arial"
    lp.font.size = Pt(9.5)
    lp.font.color.rgb = NAVY
    if note:
        nt = slide.shapes.add_textbox(Inches(x + 0.12), Inches(y + 0.72), Inches(1.4), Inches(0.15))
        np = nt.text_frame.paragraphs[0]
        np.text = note
        np.font.name = "Arial"
        np.font.size = Pt(7.5)
        np.font.color.rgb = GRAY


def add_flow_node(slide, x, y, text, color):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(1.42), Inches(0.55))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE


def add_arrow(slide, x, y):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(0.28), Inches(0.22))
    p = box.text_frame.paragraphs[0]
    p.text = "→"
    p.font.name = "Arial"
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY


def make_deck():
    REPORT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1. Title
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.18)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = BLUE
    title = slide.shapes.add_textbox(Inches(0.72), Inches(1.25), Inches(7.9), Inches(0.9))
    p = title.text_frame.paragraphs[0]
    p.text = "基于脑信号推断用户偏好的深度学习推荐实验"
    p.font.name = "Arial"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY
    sub = slide.shapes.add_textbox(Inches(0.74), Inches(2.12), Inches(7.4), Inches(0.46))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "EEG 隐式反馈  |  偏好预测  |  Top-N 推荐"
    sp.font.name = "Arial"
    sp.font.size = Pt(16)
    sp.font.color.rgb = BLUE
    add_bullets(
        slide,
        [
            "目标：用 EEG 脑电信号预测用户是否喜欢内容",
            "方法：CNN 提取脑信号特征，输出偏好概率和评分",
            "结果：将偏好概率融合物品流行度，生成推荐列表",
        ],
        0.82,
        3.05,
        6.2,
        1.4,
        15,
    )
    for i, (label, value) in enumerate([("用户", "20"), ("物品", "100"), ("样本", "2000")]):
        add_metric_card(slide, 0.82 + i * 1.88, 5.25, label, value)
    add_footer(slide, 1)

    # 2. System functions
    slide = prs.slides.add_slide(blank)
    add_title(slide, "系统功能", "从 EEG 数据到推荐结果的一体化实验流程")
    add_bullets(
        slide,
        [
            "数据模块：自动生成可复现实验 EEG 样本，包含 valence、arousal、rating 和 preference_label",
            "预处理模块：标准化、轻量滤波，并提取 delta、theta、alpha、beta、gamma 频段特征",
            "模型模块：EEGPreferenceNet 输出喜欢概率、预测评分和兴趣 embedding",
            "推荐模块：按 score = 0.7 × EEG 偏好分数 + 0.3 × 物品流行度进行 Top-N 排序",
            "结果模块：自动生成指标 CSV、推荐 CSV、图表和实验报告",
        ],
        0.7,
        1.55,
        8.4,
        3.3,
        14,
    )
    add_flow_node(slide, 0.78, 5.58, "造数", BLUE)
    add_arrow(slide, 2.25, 5.72)
    add_flow_node(slide, 2.62, 5.58, "预处理", TEAL)
    add_arrow(slide, 4.1, 5.72)
    add_flow_node(slide, 4.48, 5.58, "训练", BLUE)
    add_arrow(slide, 5.94, 5.72)
    add_flow_node(slide, 6.32, 5.58, "评估", TEAL)
    add_arrow(slide, 7.78, 5.72)
    add_flow_node(slide, 8.16, 5.58, "推荐", BLUE)
    add_footer(slide, 2)

    # 3. Technologies
    slide = prs.slides.add_slide(blank)
    add_title(slide, "实现技术", "PyTorch 深度学习模型 + 频段特征 + 推荐排序")
    add_bullets(
        slide,
        [
            "数据：模拟 DEAP 风格 EEG 情绪/偏好样本，高 valence 近似代表更强偏好",
            "特征：FFT 近似 PSD，统计五类脑电频段平均功率",
            "模型：1D CNN + BatchNorm + ReLU + Pooling + MLP",
            "输出：preference_prob、predicted_rating、embedding",
            "评估：Accuracy、Precision、Recall、F1、AUC、HitRate@10、NDCG@10、Precision@10",
        ],
        0.72,
        1.48,
        4.3,
        3.4,
        13.2,
    )
    if (FIGURE_DIR / "eeg_band_features.png").exists():
        slide.shapes.add_picture(str(FIGURE_DIR / "eeg_band_features.png"), Inches(5.42), Inches(1.65), width=Inches(3.85))
    if (FIGURE_DIR / "training_loss_curve.png").exists():
        slide.shapes.add_picture(str(FIGURE_DIR / "training_loss_curve.png"), Inches(5.42), Inches(4.42), width=Inches(3.85))
    add_footer(slide, 3)

    # 4. Results
    slide = prs.slides.add_slide(blank)
    add_title(slide, "实验结果", "分类效果较稳定，推荐结果能够命中用户偏好")
    metrics = [
        ("Accuracy", "0.820"),
        ("F1-score", "0.862"),
        ("AUC", "0.919"),
        ("HitRate@10", "1.000"),
        ("NDCG@10", "0.935"),
    ]
    for i, (label, value) in enumerate(metrics):
        add_metric_card(slide, 0.7 + i * 1.76, 1.45, label, value)
    if (FIGURE_DIR / "roc_curve.png").exists():
        slide.shapes.add_picture(str(FIGURE_DIR / "roc_curve.png"), Inches(0.72), Inches(2.75), width=Inches(3.95))
    if (FIGURE_DIR / "topn_recommendations.png").exists():
        slide.shapes.add_picture(str(FIGURE_DIR / "topn_recommendations.png"), Inches(5.08), Inches(2.75), width=Inches(4.15))
    add_footer(slide, 4)

    # 5. Summary
    slide = prs.slides.add_slide(blank)
    add_title(slide, "总结与答辩要点", "本实验完成了脑信号感知推荐的可运行闭环")
    add_bullets(
        slide,
        [
            "完成点：从 EEG 样本生成、信号处理、模型训练到 Top-N 推荐全部跑通",
            "创新点：把脑电信号作为用户隐式反馈，比点击和评分更接近即时认知状态",
            "可展示：main.py 一键运行，outputs 中保存指标、图表、推荐结果和模型文件",
            "局限性：当前为模拟数据，真实应用需要接入 DEAP/SEED 等真实 EEG 数据集",
            "改进方向：尝试 CNN-LSTM、Transformer，并融合长期行为序列提升推荐效果",
        ],
        0.85,
        1.55,
        8.2,
        3.35,
        15,
    )
    add_bullets(
        slide,
        [
            "答辩时一句话概括：本项目用深度学习从 EEG 脑信号中预测用户喜欢概率，再把它转成推荐排序分数。",
        ],
        0.85,
        5.35,
        8.2,
        0.55,
        13,
        BLUE,
    )
    add_footer(slide, 5)

    prs.save(PPTX_PATH)

    outline = """# EEG 脑信号偏好推荐实验答辩 PPT 大纲

## 第 1 页：题目与实验目标

- 用 EEG 脑电信号预测用户是否喜欢内容。
- 模型输出偏好概率和预测评分。
- 将偏好概率转成推荐排序分数。

## 第 2 页：系统功能

- 自动生成模拟 EEG 数据。
- 完成 EEG 标准化、滤波和频段特征提取。
- 训练 EEGPreferenceNet。
- 输出分类评估指标和 Top-N 推荐结果。
- 自动生成图表和实验报告。

## 第 3 页：实现技术

- 使用 PyTorch 搭建 1D CNN 模型。
- 使用 FFT/PSD 思路提取 delta、theta、alpha、beta、gamma 频段功率。
- 模型输出 preference_prob、predicted_rating 和 embedding。
- 推荐公式：score = 0.7 * eeg_preference_score + 0.3 * item_popularity_score。

## 第 4 页：实验结果

- Accuracy = 0.820。
- F1-score = 0.862。
- AUC = 0.919。
- HitRate@10 = 1.000。
- NDCG@10 = 0.935。

## 第 5 页：总结与改进

- 项目跑通了从脑信号到推荐结果的完整流程。
- 创新点是把 EEG 作为隐式反馈，比点击、评分更接近用户即时认知状态。
- 后续可以接入真实 DEAP/SEED 数据集，并尝试 CNN-LSTM 或 Transformer。
"""
    OUTLINE_PATH.write_text(outline, encoding="utf-8")
    print(f"saved pptx: {PPTX_PATH}")
    print(f"saved outline: {OUTLINE_PATH}")


if __name__ == "__main__":
    make_deck()

